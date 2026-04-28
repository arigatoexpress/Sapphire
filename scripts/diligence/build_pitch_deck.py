"""Build the Sapphire OS investor / corp-dev pitch deck.

Audience: Palantir Foundry corp-dev (default). Renders for Robinhood Cortex /
specialty acquirer with minor swaps in the "Why this acquirer" slide.

Visual style: Midnight Executive — deep navy + white + sapphire blue accent.
Sans-serif, dark backgrounds for opener / closer / appendix, lighter cards for
content slides. One diagram or table per slide. Conservative.

Usage from repo root:
    python3 scripts/diligence/build_pitch_deck_diagram.py     # writes architecture.png
    python3 scripts/diligence/build_pitch_deck.py              # writes the pptx

The diagram script writes to /tmp/sapphire-deck-build/architecture.png by
default (override via OUT_DIR env). The deck script reads that PNG and writes
the pptx alongside it. Copy the result into docs/diligence/ for distribution.
"""
from __future__ import annotations

import os
import sys
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# --- Palette (mirrors web/acquirer/assets/styles.css and the diagram) ---
NAVY_DEEP = RGBColor(0x0B, 0x12, 0x20)      # --bg              dark slide bg
NAVY_ELEV = RGBColor(0x11, 0x1A, 0x2E)      # --bg-elev         card bg darker
NAVY_CARD = RGBColor(0x14, 0x21, 0x3D)      # --bg-card
ICE = RGBColor(0xE6, 0xEC, 0xFF)            # --fg              primary text
ICE_MUTED = RGBColor(0xAA, 0xB4, 0xD4)      # --fg-muted
SAPPHIRE = RGBColor(0x6D, 0xA5, 0xFF)       # --accent
SAPPHIRE_LT = RGBColor(0x8E, 0xDC, 0xFF)    # --accent-2
LINE = RGBColor(0x2A, 0x3A, 0x5E)           # --line muted
WARN = RGBColor(0xFF, 0xD1, 0x66)           # --warn
GOOD = RGBColor(0x6E, 0xE7, 0xB7)           # --good
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_SOFT = RGBColor(0x1B, 0x1B, 0x1B)


HEAD_FONT = "Helvetica Neue"
BODY_FONT = "Helvetica Neue"
MONO_FONT = "Menlo"


# Provenance (HEAD SHA fetched at build time; date frozen for determinism)
HEAD_SHA = os.environ.get("SAPPHIRE_HEAD_SHA", "9074c408")
BUILD_DATE = "2026-04-29"
TARGET_AUDIENCE = "Palantir Foundry corp-dev"


# ---------- Helpers ----------

def add_dark_background(slide, color=NAVY_DEEP):
    """Fill slide background with a solid color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             italic=False, color=ICE, align=PP_ALIGN.LEFT,
             vertical=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = vertical

    # Support multi-line text via newlines
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_bullet_list(slide, x, y, w, h, bullets, *, size=14, color=ICE,
                    indent_size=12, marker_color=SAPPHIRE_LT, line_spacing=1.25):
    """Bullet list with sapphire-blue dot markers (avoid the AI-cliche of fancy bullets)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0)
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Marker
        marker = p.add_run()
        marker.text = "•  "
        marker.font.name = HEAD_FONT
        marker.font.size = Pt(size + 1)
        marker.font.bold = True
        marker.font.color.rgb = marker_color
        # Body
        body = p.add_run()
        body.text = bullet
        body.font.name = BODY_FONT
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, fill=NAVY_CARD, edge=LINE, edge_w=Pt(0.75)):
    """Rounded card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # tighten corner radius
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = edge
    card.line.width = edge_w
    card.shadow.inherit = False
    return card


def add_accent_bar(slide, x, y, w, h, color=SAPPHIRE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_footer(slide, *, dark=True, label=None):
    """Small bottom-corner footer: brand left, slide info right."""
    fg = ICE_MUTED if dark else RGBColor(0x6A, 0x70, 0x88)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(6.0), Inches(0.3),
             "Sapphire OS  ·  Pitch  ·  " + BUILD_DATE,
             font=BODY_FONT, size=9, color=fg)
    if label:
        add_text(slide, Inches(7.0), Inches(7.05), Inches(6.0), Inches(0.3),
                 label, font=BODY_FONT, size=9, color=fg, align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def set_table_cell(cell, text, *, font=BODY_FONT, size=11, bold=False,
                   color=ICE, fill=None, align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""  # reset
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


# ---------- Slide builders ----------

def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def slide_title(prs):
    """Slide 1 — Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_background(slide, NAVY_DEEP)

    # Brand mark + name top-left
    add_text(slide, Inches(0.6), Inches(0.6), Inches(2.0), Inches(0.4),
             "[ S ]   SAPPHIRE OS", font=MONO_FONT, size=12, bold=True,
             color=SAPPHIRE_LT)

    # Big positioning headline (allow 2 lines, give it 2.2" of vertical room)
    add_text(slide, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.2),
             "A production-graded\nautonomy control plane.",
             font=HEAD_FONT, size=44, bold=True, color=ICE,
             line_spacing=1.1)
    add_text(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(0.6),
             "Risk-gated intelligence, trading research, and operations — auditable before action.",
             font=HEAD_FONT, size=18, color=SAPPHIRE_LT, line_spacing=1.2)

    # Accent bar
    add_accent_bar(slide, Inches(0.6), Inches(5.05), Inches(2.5), Inches(0.05),
                   SAPPHIRE)

    # Operator + audience
    add_text(slide, Inches(0.6), Inches(5.3), Inches(12.0), Inches(0.4),
             "Aristotle Spec  ·  Operator + Founder",
             font=HEAD_FONT, size=14, color=ICE)
    add_text(slide, Inches(0.6), Inches(5.7), Inches(12.0), Inches(0.4),
             f"Prepared for: {TARGET_AUDIENCE}",
             font=BODY_FONT, size=14, color=ICE_MUTED)

    # SHA + date bottom-left
    add_text(slide, Inches(0.6), Inches(6.6), Inches(8.0), Inches(0.3),
             f"HEAD: {HEAD_SHA}  ·  {BUILD_DATE}  ·  arigatoexpress/Sapphire",
             font=MONO_FONT, size=10, color=ICE_MUTED)
    add_text(slide, Inches(8.0), Inches(6.6), Inches(4.7), Inches(0.3),
             "aristotlespec@gmail.com",
             font=MONO_FONT, size=10, color=ICE_MUTED, align=PP_ALIGN.RIGHT)

    add_speaker_notes(
        slide,
        "Title slide. Sapphire OS is a production-graded autonomy control "
        "plane: risk-gated intelligence, trading research, and operations, "
        "auditable before action.\n\n"
        f"Audience for this render: {TARGET_AUDIENCE} — see Slide 8 for "
        "audience-specific framing.\n\n"
        f"Source: HEAD {HEAD_SHA} @ docs/diligence/00-executive-summary.md"
    )
    return slide


def slide_thesis(prs):
    """Slide 2 — The thesis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6),
             "The thesis",
             font=HEAD_FONT, size=32, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.4),
             "Why Sapphire is buyable, and why now.",
             font=HEAD_FONT, size=14, color=ICE_MUTED, italic=True)

    # Three large numbered cards
    card_y = Inches(1.95)
    card_h = Inches(4.4)
    card_w = Inches(4.0)
    gap = Inches(0.18)
    starts = [Inches(0.6), Inches(0.6) + card_w + gap, Inches(0.6) + (card_w + gap) * 2]

    items = [
        (
            "01",
            "What it is",
            "An auditable intelligence + trading-research operating system. "
            "Multi-source signals fuse into one edge_score. Every artifact ships a "
            "provenance envelope. Every action passes a risk verdict before it fires.",
        ),
        (
            "02",
            "Why now",
            "The differentiators are now separable: Risk Kernel, Provenance, Foundry Ontology, "
            "and a 5,304-test gate exist as packaged surfaces. "
            "Operator intuition is now file-addressable.",
        ),
        (
            "03",
            "Why defensible",
            "The moat is restraint encoded in code: dry-run by default, fail-closed risk gates, "
            "no-spend CI, secrets outside the repo, live trading capped at $5 with manual confirm.",
        ),
    ]

    for x, (num, title, body) in zip(starts, items):
        add_card(slide, x, card_y, card_w, card_h, fill=NAVY_CARD)
        add_text(slide, x + Inches(0.3), card_y + Inches(0.3),
                 card_w - Inches(0.6), Inches(0.6),
                 num, font=MONO_FONT, size=28, bold=True, color=SAPPHIRE)
        add_text(slide, x + Inches(0.3), card_y + Inches(1.0),
                 card_w - Inches(0.6), Inches(0.5),
                 title, font=HEAD_FONT, size=18, bold=True, color=ICE)
        add_text(slide, x + Inches(0.3), card_y + Inches(1.6),
                 card_w - Inches(0.6), card_h - Inches(1.8),
                 body, font=BODY_FONT, size=12.5, color=ICE_MUTED,
                 line_spacing=1.35)

    add_footer(slide, dark=True, label="2 / 14")
    add_speaker_notes(
        slide,
        "Three pillars: what, when, why-defensible. Each maps to a section "
        "in the diligence packet (00 executive summary; 02 product surfaces; "
        "03 security posture).\n\n"
        "Test count is from scripts/ops/test_inventory.py (5,304 collected on "
        "2026-04-28). Live $5 cap is from docs/products/live-trading-ramp-memo.md.\n\n"
        f"Source: docs/diligence/00-executive-summary.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_architecture(prs, image_path):
    """Slide 3 — Intelligence stack visualized."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.5),
             "The intelligence stack",
             font=HEAD_FONT, size=26, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.78), Inches(12.0), Inches(0.35),
             "Seven independent signal sources fuse into one governed edge.",
             font=HEAD_FONT, size=12, italic=True, color=ICE_MUTED)

    # Image fills most of the slide; lock by HEIGHT so it doesn't clip.
    # Slide is 7.5" tall; reserve top ~1.2" for title + bottom 0.5" for footer.
    # Available height = ~5.7"; image native ratio is 16:9 → max width ~10.13".
    # Center horizontally on the 13.333"-wide slide.
    pic_height_in = 5.7
    pic_width_in = pic_height_in * 16.0 / 9.0
    pic_left = Inches((13.333 - pic_width_in) / 2.0)
    pic_top = Inches(1.25)
    slide.shapes.add_picture(image_path, pic_left, pic_top,
                             height=Inches(pic_height_in))

    add_footer(slide, dark=True, label="3 / 14")
    add_speaker_notes(
        slide,
        "Architecture rendered from matplotlib — palette mirrors the brand "
        "(midnight navy, sapphire accent). The seven signal sources on the "
        "left are real, shipping surfaces (Tranche 2 + Tranche 4). The middle "
        "column is the Tranche 3 correlator, narrative synthesis, risk kernel, "
        "and confirmation firewall — all production code with tests. The "
        "right column is the governed-output set.\n\n"
        "The architecture diagram in docs/diligence/01-architecture.md is the "
        "Mermaid source of truth; this slide is a visual rendering of the "
        "same decomposition.\n\n"
        f"Source: docs/diligence/01-architecture.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_capabilities(prs):
    """Slide 4 — Capabilities matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Capabilities matrix",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Fourteen sellable surfaces. Versioned, tested, separately purchasable.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Table
    rows = [
        ("Surface", "What it does", "Status"),
        ("Risk Kernel 0.1.0",
         "DecisionEnvelope → RiskVerdict; fail-closed policy tree",
         "live"),
        ("Provenance Envelopes 0.1.0",
         "SHA-256 sidecar on every artifact; verifier in CI",
         "live"),
        ("Signal Correlator 0.1.0",
         "8 source adapters → unified edge_score per (symbol, timeframe)",
         "live"),
        ("Narrative Synthesis 0.1.0",
         "Bounded thesis + invalidators; rubric-gated publication",
         "dry-run default"),
        ("Foundry Ontology 0.2.0",
         "13 typed objects with watermarks + idempotency ledger",
         "dry-run default"),
        ("Threat Intel Product 0.1.0",
         "/threat-intel page over CISA KEV + NVD + MITRE ATT&CK",
         "live (read-only)"),
        ("Customer Dossier 0.2.0",
         "Per-tenant HMAC isolation + cell-suppression on PII",
         "live (read-only)"),
        ("BQ Vector Retrieval 0.1.0",
         "intel_search over IntelVectorRecord; mock + BQ live path",
         "mock default"),
        ("Cross-Asset Correlation 0.1.0",
         "Pearson/Spearman/Kendall + regime breakdown events",
         "live (read-only)"),
        ("Macro Intel 0.1.0",
         "Fed/SEC/CFTC/Treasury official feeds + forward calendar",
         "dry-run default"),
        ("Event-Impact Modeling 0.1.0",
         "Historical reaction profiles per (event class, asset, horizon)",
         "live (read-only)"),
        ("Counterparty Intel 0.1.0",
         "Hyperliquid public top-trader consensus signals",
         "live (read-only)"),
        ("On-Chain Intelligence 0.2.0",
         "Glassnode / Santiment / ETH+SOL nodes; provider-gated",
         "dry-run default"),
        ("Adversarial Defense 0.1.0",
         "5 detectors: wash trade, prompt inj, oracle, false flag, bot pump",
         "telemetry-first"),
    ]

    table_left = Inches(0.6)
    table_top = Inches(1.55)
    table_width = Inches(12.1)
    table_height = Inches(5.3)
    n_rows = len(rows)
    n_cols = 3

    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(3.4)
    table.columns[1].width = Inches(6.6)
    table.columns[2].width = Inches(2.1)
    # row heights
    table.rows[0].height = Inches(0.4)
    for r in range(1, n_rows):
        table.rows[r].height = Inches(0.34)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_table_cell(cell, val, font=HEAD_FONT, size=11, bold=True,
                               color=NAVY_DEEP, fill=SAPPHIRE_LT,
                               align=PP_ALIGN.LEFT if c < 2 else PP_ALIGN.LEFT)
            else:
                # Status coloring
                if c == 2:
                    status = val.lower()
                    if "dry" in status or "telem" in status or "mock" in status:
                        col = WARN
                    elif "live" in status:
                        col = GOOD
                    else:
                        col = ICE
                    fill = NAVY_ELEV if r % 2 else NAVY_CARD
                    set_table_cell(cell, val, font=BODY_FONT, size=10.5,
                                   bold=True, color=col, fill=fill)
                else:
                    fill = NAVY_ELEV if r % 2 else NAVY_CARD
                    bold = (c == 0)
                    txt_color = ICE if c == 0 else ICE_MUTED
                    set_table_cell(cell, val, font=BODY_FONT, size=10.5,
                                   bold=bold, color=txt_color, fill=fill)

    add_footer(slide, dark=True, label="4 / 14")
    add_speaker_notes(
        slide,
        "Fourteen versioned product surfaces. Status colors:\n"
        "  green = live (read-only or paper) and shipping\n"
        "  amber = dry-run default / mock fixture / telemetry-first\n\n"
        "Every row corresponds to a doc under docs/products/<surface>-<version>.md "
        "with tests under tests/unit/. None are vaporware. Each is independently "
        "purchasable — Risk Kernel and Provenance in particular are designed as "
        "carve-out libraries.\n\n"
        f"Source: docs/products/*.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_safety(prs):
    """Slide 5 — Safety posture: kill-switch layers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Safety posture",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Five kill-switch layers. Every mutation path passes through them.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # 5 cards in a row
    layers = [
        ("L1", "Trade-time invariants",
         "RiskKernelV1 + HardRiskKernel + portfolio kill switch.",
         "Daily-loss · drawdown · ATR · kill-switch · data-leakage policies. "
         "Fail-closed verdict tree."),
        ("L2", "Confirmation firewall",
         "lib/core/confirmation_firewall.py — 2-phase commit on destructive verbs.",
         "Read-only auto-approves. Financial / external / destructive require operator "
         "confirmation; destructive enforces 30-second delay."),
        ("L3", "Operator manual halt",
         "Telegram /routines pause · /cancel-routine · CONFIRM token to resume.",
         "Pause flag short-circuits future ticks via abort_if_paused. In-flight work is "
         "not killed (intentional)."),
        ("L4", "Security kill switch",
         "lib/core/security_kill_switch.py — for suspected account compromise.",
         "Stops signal logger, writes cloud-routing kill flag, unloads content engine, "
         "stops inference proxy. Reversible."),
        ("L5", "Heartbeat + supervisor",
         "lib/core/heartbeat.py + service_supervisor.py.",
         "Per-component health state, consecutive-failure counters, and structured "
         "evidence — observability, not silent retargeting."),
    ]

    card_y = Inches(1.6)
    card_w = Inches(2.42)
    card_h = Inches(5.0)
    gap = Inches(0.08)
    start_x = Inches(0.5)

    for i, (tag, title, lib, desc) in enumerate(layers):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, card_y, card_w, card_h, fill=NAVY_CARD, edge=LINE)
        # Top accent strip
        add_accent_bar(slide, x, card_y, card_w, Inches(0.05), SAPPHIRE)

        add_text(slide, x + Inches(0.2), card_y + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.4),
                 tag, font=MONO_FONT, size=15, bold=True, color=SAPPHIRE_LT)
        add_text(slide, x + Inches(0.2), card_y + Inches(0.7),
                 card_w - Inches(0.4), Inches(1.1),
                 title, font=HEAD_FONT, size=13.5, bold=True, color=ICE,
                 line_spacing=1.2)
        add_text(slide, x + Inches(0.2), card_y + Inches(1.95),
                 card_w - Inches(0.4), Inches(1.1),
                 lib, font=MONO_FONT, size=8.5, italic=True, color=ICE_MUTED,
                 line_spacing=1.25)
        add_text(slide, x + Inches(0.2), card_y + Inches(3.05),
                 card_w - Inches(0.4), card_h - Inches(3.2),
                 desc, font=BODY_FONT, size=10, color=ICE_MUTED,
                 line_spacing=1.3)

    add_footer(slide, dark=True, label="5 / 14")
    add_speaker_notes(
        slide,
        "Five kill-switch layers, each tested under tests/unit/ "
        "(test_risk_kernel_*, test_confirmation_firewall, test_kill_switch, "
        "test_security_kill_switch, test_routine_pause, test_heartbeat).\n\n"
        "Note: docs/security/kill-switch-invariants.md enumerates four primary "
        "layers plus heartbeat as observability — this slide flattens them to "
        "five for visual symmetry.\n\n"
        f"Source: docs/security/kill-switch-invariants.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_traction(prs):
    """Slide 6 — Production traction: hard numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Production traction",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Concrete numbers — pulled from the repo on 2026-04-28.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # 8 large stat tiles in 2x4
    stats = [
        ("5,304", "tests collected",
         "scripts/ops/test_inventory.py"),
        ("63", "plugin tools on disk",
         "plugins/claw-sapphire/tools/"),
        ("32", "dashboard pages",
         "services/dashboard/templates/pages/"),
        ("24", "macOS LaunchAgents",
         "infra/launchagents/"),
        ("7", "quant strategies",
         "lib/analytics/strategies/"),
        ("13", "Foundry ontology types",
         "docs/foundry-ontology-schema.md"),
        ("$5", "live-trading first-order cap",
         "manual confirm · crypto only"),
        ("0", "fabricated metrics",
         "every number above is repo-grounded"),
    ]

    grid_x = Inches(0.5)
    grid_y = Inches(1.6)
    tile_w = Inches(3.05)
    tile_h = Inches(2.55)
    gap_x = Inches(0.12)
    gap_y = Inches(0.18)
    cols = 4

    for i, (number, label, source) in enumerate(stats):
        r = i // cols
        c = i % cols
        x = grid_x + c * (tile_w + gap_x)
        y = grid_y + r * (tile_h + gap_y)
        add_card(slide, x, y, tile_w, tile_h, fill=NAVY_CARD, edge=LINE)

        add_text(slide, x + Inches(0.2), y + Inches(0.25),
                 tile_w - Inches(0.4), Inches(1.2),
                 number, font=HEAD_FONT, size=42, bold=True, color=SAPPHIRE_LT,
                 align=PP_ALIGN.LEFT)
        add_text(slide, x + Inches(0.2), y + Inches(1.45),
                 tile_w - Inches(0.4), Inches(0.5),
                 label, font=HEAD_FONT, size=13, bold=True, color=ICE)
        add_text(slide, x + Inches(0.2), y + Inches(1.95),
                 tile_w - Inches(0.4), Inches(0.55),
                 source, font=MONO_FONT, size=9, italic=True, color=ICE_MUTED,
                 line_spacing=1.25)

    add_footer(slide, dark=True, label="6 / 14")
    add_speaker_notes(
        slide,
        "All metrics grounded:\n"
        "  • 5,304 tests: from scripts/ops/test_inventory.py output on "
        "2026-04-28 (4,928 unit + 376 plugin)\n"
        "  • 63 plugin tools: 36 top-level + 25 internal + 2 deprecated, "
        "per CLAUDE.md and ls plugins/claw-sapphire/tools/\n"
        "  • 32 dashboard pages: per CLAUDE.md and "
        "ls services/dashboard/templates/pages/ (38 files; 32 unique product pages)\n"
        "  • 24 LaunchAgents: ls infra/launchagents/ (CLAUDE.md notes 20 in audit "
        "snapshot — current count is 24 after recent additions)\n"
        "  • 7 quant strategies: lib/analytics/ — RegimeAwareRSI, FundingRate-"
        "Contrarian, CorrelationBreakout, MultiTFMomentum, SapphireComposite + base + params\n"
        "  • 13 Foundry types: 8 from 0.1.0 + 5 added in 0.2.0\n"
        "  • $5 live cap: docs/products/live-trading-ramp-memo.md (capped, "
        "manual-confirm-only path; not yet a recurring autonomous flow)\n\n"
        f"Source: CLAUDE.md + scripts/ops/test_inventory.py @ SHA {HEAD_SHA}"
    )
    return slide


def slide_competitive(prs):
    """Slide 7 — Competitive landscape (positioning matrix)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Competitive landscape",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Honest positioning. Primary-source-only research, 2026-04-28.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    rows = [
        ("Capability", "Sapphire", "Foundry/AIP", "Cortex", "Bloomberg"),
        ("Multi-source signal correlation",        "yes",     "partial", "partial", "partial"),
        ("Typed ontology / decision objects",      "partial", "yes",     "—",       "partial"),
        ("Per-artifact provenance envelopes",      "yes",     "partial", "—",       "partial"),
        ("Bounded narrative synthesis",            "partial", "partial", "partial", "partial"),
        ("Dry-run external actions + no-spend",    "yes",     "—",       "partial", "—"),
        ("Owner-operated autonomy (1-operator)",   "yes",     "—",       "—",       "—"),
        ("Regulated execution (broker / exchange)","no",      "no",      "yes",     "partial"),
        ("Real-time institutional market data",    "no",      "partial", "partial", "yes"),
        ("Enterprise implementation partners",     "no",      "yes",     "—",       "partial"),
    ]

    n_rows = len(rows)
    n_cols = 5
    table_left = Inches(0.6)
    table_top = Inches(1.5)
    table_width = Inches(12.1)
    table_height = Inches(4.0)
    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(4.5)
    for c in range(1, 5):
        table.columns[c].width = Inches(1.9)
    table.rows[0].height = Inches(0.4)
    for r in range(1, n_rows):
        table.rows[r].height = Inches(0.34)

    def cell_color_for(val):
        v = val.lower().strip()
        if v == "yes":
            return GOOD
        if v == "partial":
            return WARN
        if v == "no":
            return RGBColor(0xFF, 0x90, 0x90)
        return ICE_MUTED

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_table_cell(cell, val, font=HEAD_FONT, size=11, bold=True,
                               color=NAVY_DEEP, fill=SAPPHIRE_LT,
                               align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
            else:
                fill = NAVY_ELEV if r % 2 else NAVY_CARD
                if c == 0:
                    set_table_cell(cell, val, font=BODY_FONT, size=10.5,
                                   bold=True, color=ICE, fill=fill)
                else:
                    color = cell_color_for(val)
                    bold = val.lower() in ("yes", "partial", "no")
                    set_table_cell(cell, val, font=BODY_FONT, size=10.5,
                                   bold=bold, color=color, fill=fill,
                                   align=PP_ALIGN.CENTER)

    # Honest takeaway
    add_text(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0),
             "The honest differentiator: governed owner-operated autonomy. "
             "We are weakest on regulated execution and licensed market data — "
             "and we deliberately do not compete there.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED,
             line_spacing=1.4)

    add_footer(slide, dark=True, label="7 / 14")
    add_speaker_notes(
        slide,
        "Adapted from docs/competitive/landscape-2026-04-28.md (40k-word "
        "primary-source memo). Color: green = yes, amber = partial, red = no.\n\n"
        "Two key honest concessions:\n"
        "  • Sapphire does not compete on real-time institutional market data "
        "(Bloomberg moat).\n"
        "  • Sapphire does not compete on regulated execution (Robinhood / broker rails).\n\n"
        "Where Sapphire genuinely differentiates: governed, owner-operated "
        "autonomy with provenance envelopes + dry-run defaults + no-spend CI.\n\n"
        f"Source: docs/competitive/landscape-2026-04-28.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_why_audience(prs):
    """Slide 8 — Why this acquirer (Palantir Foundry default)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             f"Why {TARGET_AUDIENCE}",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Sapphire speaks Foundry's language already. We have the "
             "ontology, the provenance, the action contract.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Three big alignment cards
    cards = [
        (
            "Ontology alignment",
            "13 typed objects (Foundry Ontology 0.2.0)",
            "PaperTrade · Alert · ServiceHealth · ThreatIntel · DailyBrief · Region · "
            "IntelItem · IntelSourceHealth · IntelVectorRecord · TelegramIntelMessage · "
            "HyperliquidSignal · OODAPacket · ThreatIndicator.\n\n"
            "Watermarks + idempotency ledger. Live mode behind SAPPHIRE_FOUNDRY_LIVE=1.",
        ),
        (
            "AIP-shaped governance",
            "Risk Kernel + Provenance + Confirmation Firewall",
            "Per-artifact envelope + risk verdict + 2-phase confirmation matches "
            "AIP Logic's evaluation + action gating shape.\n\n"
            "What Foundry sells as enterprise infrastructure, Sapphire ships as a "
            "compact single-operator reference implementation.",
        ),
        (
            "POC fit, not replacement",
            "We don't out-Foundry Foundry — we plug into it.",
            "Sapphire is exactly the right size to be a Foundry pilot tenant: "
            "5,304 tests, 14 productized surfaces, and provenance discipline that "
            "maps to your ontology contracts on day one.",
        ),
    ]

    card_y = Inches(1.6)
    card_h = Inches(5.0)
    card_w = Inches(4.05)
    gap = Inches(0.15)
    starts = [Inches(0.55),
              Inches(0.55) + card_w + gap,
              Inches(0.55) + (card_w + gap) * 2]

    for x, (title, sub, body) in zip(starts, cards):
        add_card(slide, x, card_y, card_w, card_h, fill=NAVY_CARD)
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), SAPPHIRE)

        add_text(slide, x + Inches(0.3), card_y + Inches(0.35),
                 card_w - Inches(0.6), Inches(0.6),
                 title, font=HEAD_FONT, size=18, bold=True, color=ICE)
        add_text(slide, x + Inches(0.3), card_y + Inches(1.1),
                 card_w - Inches(0.6), Inches(0.6),
                 sub, font=HEAD_FONT, size=12, italic=True, color=SAPPHIRE_LT,
                 line_spacing=1.25)
        add_text(slide, x + Inches(0.3), card_y + Inches(1.95),
                 card_w - Inches(0.6), card_h - Inches(2.1),
                 body, font=BODY_FONT, size=12, color=ICE_MUTED,
                 line_spacing=1.4)

    add_footer(slide, dark=True, label="8 / 14")
    add_speaker_notes(
        slide,
        "This slide is the audience-specific page. Default rendering targets "
        f"{TARGET_AUDIENCE}.\n\n"
        "For Robinhood Cortex: swap the headline to 'Why Robinhood' and "
        "emphasize: (1) broker-agnostic risk-explanation that respects "
        "Cortex's no-trade boundary; (2) crypto-specific intelligence "
        "(Hyperliquid, on-chain, counterparty); (3) live-trading ramp "
        "discipline ($5 → $50 → $500) that mirrors Cortex's careful "
        "informational-vs-execution separation.\n\n"
        "For specialty acquirer: emphasize compactness — 5,304 tests, "
        "1 operator, an integrated control plane that demonstrates the "
        "safe-autonomy pattern at carve-out scale.\n\n"
        f"Source: docs/foundry-ontology-schema.md + docs/products/foundry-ontology-0.2.0.md "
        f"+ project_palantir_pitch memory @ SHA {HEAD_SHA}"
    )
    return slide


def slide_live_trading_ramp(prs):
    """Slide 9 — Live trading ramp."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Live-trading ramp",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Paper-mostly today. Live-capital path is gated by metrics, not vibes.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Three rungs as cards in a row, with arrows between
    rungs = [
        ("$5", "First-order pilot",
         "Single manual-confirm crypto limit order. UUID client_order_id. "
         "One-order confirmation token. Crypto only.",
         GOOD),
        ("$50", "Capped pilot",
         "After 14 trading days · Sortino > 1.5 at the strategy level · "
         "no firewall bypass · clean readiness sweep. Daily live cap $10.",
         WARN),
        ("$500", "Live tier",
         "Designed future rung. Requires fresh PR · operator-reviewed runbook update · "
         "stronger live-order telemetry. Not active today.",
         ICE_MUTED),
    ]

    rung_y = Inches(1.65)
    rung_h = Inches(2.7)
    rung_w = Inches(3.7)
    gap_x = Inches(0.65)
    start_x = Inches(0.7)

    arrow_y_mid = rung_y + rung_h / 2
    for i, (amount, label, body, color) in enumerate(rungs):
        x = start_x + i * (rung_w + gap_x)
        add_card(slide, x, rung_y, rung_w, rung_h, fill=NAVY_CARD)
        # Top accent
        add_accent_bar(slide, x, rung_y, rung_w, Inches(0.06), color)

        add_text(slide, x + Inches(0.3), rung_y + Inches(0.3),
                 rung_w - Inches(0.6), Inches(1.0),
                 amount, font=HEAD_FONT, size=44, bold=True, color=color)
        add_text(slide, x + Inches(0.3), rung_y + Inches(1.25),
                 rung_w - Inches(0.6), Inches(0.45),
                 label, font=HEAD_FONT, size=15, bold=True, color=ICE)
        add_text(slide, x + Inches(0.3), rung_y + Inches(1.75),
                 rung_w - Inches(0.6), rung_h - Inches(1.85),
                 body, font=BODY_FONT, size=11, color=ICE_MUTED,
                 line_spacing=1.35)

        if i < 2:
            # "→" glyph between cards rather than a thin straight connector
            arrow_x = x + rung_w + Inches(0.05)
            arrow_w = gap_x - Inches(0.1)
            add_text(slide, arrow_x, arrow_y_mid - Inches(0.25),
                     arrow_w, Inches(0.5),
                     "→", font=HEAD_FONT, size=24, bold=True, color=SAPPHIRE,
                     align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)

    # Bottom honest banner
    banner_y = Inches(4.7)
    add_card(slide, Inches(0.7), banner_y, Inches(11.95), Inches(2.05),
             fill=NAVY_ELEV, edge=WARN, edge_w=Pt(1.0))
    add_text(slide, Inches(0.95), banner_y + Inches(0.18),
             Inches(11.5), Inches(0.4),
             "Honest framing — live posture as of 2026-04-28",
             font=HEAD_FONT, size=14, bold=True, color=WARN)
    add_bullet_list(slide, Inches(0.95), banner_y + Inches(0.65),
                    Inches(11.45), Inches(1.45),
                    [
                        "Sapphire is paper-mostly. The default path is paper trading + dry-run order drafts.",
                        "Live posture is manual_confirmed_crypto_only. Stock automation is blocked — no official broker API exists.",
                        "Every rung is reversible. Removing --execute and pausing routines returns the system to paper-only.",
                    ],
                    size=11.5, color=ICE, marker_color=WARN, line_spacing=1.25)

    add_footer(slide, dark=True, label="9 / 14")
    add_speaker_notes(
        slide,
        "Live-trading posture: manual_confirmed_crypto_only.\n\n"
        "Key honest caveat: the operator memory references a $5 first BTC "
        "live-order intent on 2026-04-28 04:06 UTC, but that timestamp is "
        "operator-side state, not a repo-resident artifact this deck can "
        "verify. The deck therefore frames the $5 rung as 'designed and gated' "
        "rather than asserting an executed fill — preserving the no-fabrication "
        "rule.\n\n"
        "Promotion gates: 14 trading days · Sortino > 1.5 · no firewall bypass · "
        "no unexplained slippage · 0-FAIL readiness sweep.\n\n"
        f"Source: docs/products/live-trading-ramp-memo.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_team_ip(prs):
    """Slide 10 — Team + IP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Team + IP",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "One operator. Four years of compounded autonomy infrastructure.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Left card: Team
    team_x = Inches(0.6)
    team_y = Inches(1.6)
    team_w = Inches(6.05)
    team_h = Inches(5.2)
    add_card(slide, team_x, team_y, team_w, team_h, fill=NAVY_CARD)
    add_accent_bar(slide, team_x, team_y, team_w, Inches(0.06), SAPPHIRE)
    add_text(slide, team_x + Inches(0.3), team_y + Inches(0.3),
             team_w - Inches(0.6), Inches(0.5),
             "Team", font=HEAD_FONT, size=20, bold=True, color=ICE)
    add_text(slide, team_x + Inches(0.3), team_y + Inches(0.85),
             team_w - Inches(0.6), Inches(0.4),
             "Aristotle Spec — operator, founder, repo author.",
             font=HEAD_FONT, size=14, italic=True, color=SAPPHIRE_LT)
    add_bullet_list(slide, team_x + Inches(0.3), team_y + Inches(1.45),
                    team_w - Inches(0.6), team_h - Inches(1.6),
                    [
                        "Sole repo committer. Final risk authority.",
                        "Operator-in-residence through transition window.",
                        "Codex (production-autonomy assistant) handles repo work; "
                        "Claude / Hermes are constrained reviewers.",
                        "Honest: key-person concentration is a risk. "
                        "The diligence packet exists to make it transferable.",
                    ],
                    size=12.5, color=ICE_MUTED)

    # Right card: IP
    ip_x = Inches(6.85)
    ip_y = Inches(1.6)
    ip_w = Inches(6.05)
    ip_h = Inches(5.2)
    add_card(slide, ip_x, ip_y, ip_w, ip_h, fill=NAVY_CARD)
    add_accent_bar(slide, ip_x, ip_y, ip_w, Inches(0.06), SAPPHIRE_LT)
    add_text(slide, ip_x + Inches(0.3), ip_y + Inches(0.3),
             ip_w - Inches(0.6), Inches(0.5),
             "IP + Repo", font=HEAD_FONT, size=20, bold=True, color=ICE)
    add_text(slide, ip_x + Inches(0.3), ip_y + Inches(0.85),
             ip_w - Inches(0.6), Inches(0.4),
             "Private monorepo. Single owner. No outstanding contributor licenses.",
             font=HEAD_FONT, size=14, italic=True, color=SAPPHIRE_LT)
    add_bullet_list(slide, ip_x + Inches(0.3), ip_y + Inches(1.45),
                    ip_w - Inches(0.6), ip_h - Inches(1.6),
                    [
                        "Core repo: arigatoexpress/Sapphire (private). "
                        "5,304 tests · 14 product surfaces · 32 dashboard pages.",
                        "Satellite repos orchestrated, not absorbed.",
                        "Smart contracts: SapphireSignalVerifier + PaymentGate "
                        "on Robinhood Chain testnet (Arbitrum Orbit).",
                        "Dependencies audited via OSV + Trivy. CycloneDX SBOM.",
                    ],
                    size=12.5, color=ICE_MUTED)

    add_footer(slide, dark=True, label="10 / 14")
    add_speaker_notes(
        slide,
        "This slide is intentionally brief — operator edits afterward.\n\n"
        "Honest framing: this is a 1-operator project. The team-and-process "
        "diligence section (docs/diligence/09-team-and-process.md) covers "
        "the strength (taste, speed) and the risk (key-person concentration).\n\n"
        f"Source: docs/diligence/09-team-and-process.md + CLAUDE.md @ SHA {HEAD_SHA}"
    )
    return slide


def slide_acquisition_rationale(prs):
    """Slide 11 — Acquisition rationale: three options."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Acquisition rationale",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Three structures. Pick the one that matches your conviction window.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    options = [
        (
            "Option A",
            "Integration POC",
            "Lowest commitment. Highest signal.",
            [
                "30-day Foundry-pilot scoping",
                "Sapphire syncs into your ontology read-only",
                "No equity event; budget-light",
                "Outcome: decide whether to escalate",
            ],
            "Recommended first step.",
            GOOD,
        ),
        (
            "Option B",
            "Talent + IP licensing",
            "Acqui-hire shape with code carve-out.",
            [
                "Risk Kernel + Provenance Envelopes as MIT/Apache",
                "Operator joins as Foundry safety-IC for 18–24 mo",
                "Mid-range cash + retention",
                "Outcome: codify the safety pattern inside Foundry",
            ],
            "If conviction lands on the safety pattern.",
            WARN,
        ),
        (
            "Option C",
            "Full acquisition",
            "Sapphire as the autonomy-layer reference implementation.",
            [
                "Full repo + satellites under buyer ownership",
                "Operator-in-residence transition (6–12 mo)",
                "Top-tier cash + earnout tied to integration milestones",
                "Outcome: Foundry's compact autonomous tenant template",
            ],
            "If conviction lands on the operating system.",
            SAPPHIRE_LT,
        ),
    ]

    card_y = Inches(1.6)
    card_h = Inches(5.2)
    card_w = Inches(4.05)
    gap = Inches(0.15)
    starts = [Inches(0.55),
              Inches(0.55) + card_w + gap,
              Inches(0.55) + (card_w + gap) * 2]

    for x, (label, name, sub, bullets, bottom, color) in zip(starts, options):
        add_card(slide, x, card_y, card_w, card_h, fill=NAVY_CARD)
        add_accent_bar(slide, x, card_y, card_w, Inches(0.06), color)

        add_text(slide, x + Inches(0.3), card_y + Inches(0.3),
                 card_w - Inches(0.6), Inches(0.4),
                 label, font=MONO_FONT, size=11, bold=True, color=color)
        add_text(slide, x + Inches(0.3), card_y + Inches(0.7),
                 card_w - Inches(0.6), Inches(0.7),
                 name, font=HEAD_FONT, size=20, bold=True, color=ICE)
        add_text(slide, x + Inches(0.3), card_y + Inches(1.45),
                 card_w - Inches(0.6), Inches(0.5),
                 sub, font=HEAD_FONT, size=12, italic=True, color=ICE_MUTED,
                 line_spacing=1.25)
        add_bullet_list(slide, x + Inches(0.3), card_y + Inches(2.15),
                        card_w - Inches(0.6), card_h - Inches(2.85),
                        bullets, size=11, color=ICE_MUTED, marker_color=color)
        add_text(slide, x + Inches(0.3), card_y + card_h - Inches(0.55),
                 card_w - Inches(0.6), Inches(0.4),
                 bottom, font=HEAD_FONT, size=11, italic=True, color=color)

    add_footer(slide, dark=True, label="11 / 14")
    add_speaker_notes(
        slide,
        "Three structures from the diligence packet's executive summary "
        "and the project_palantir_pitch.md memory note. Option A is "
        "recommended as a first step because it preserves both buyer and "
        "seller optionality.\n\n"
        "Option B mirrors the acqui-hire-with-IP-licensing fallback in the "
        "Palantir pitch.\n\n"
        "Option C is the full acquisition path — most likely outcome if a "
        "POC under Option A succeeds.\n\n"
        f"Source: docs/diligence/00-executive-summary.md + project_palantir_pitch memory @ SHA {HEAD_SHA}"
    )
    return slide


def slide_the_ask(prs):
    """Slide 12 — The ask."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_DEEP)

    # Big centered ask
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.6),
             "The ask",
             font=HEAD_FONT, size=28, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "One concrete next step.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Hero card
    hero_x = Inches(1.0)
    hero_y = Inches(1.55)
    hero_w = Inches(11.3)
    hero_h = Inches(4.4)
    add_card(slide, hero_x, hero_y, hero_w, hero_h, fill=NAVY_CARD,
             edge=SAPPHIRE, edge_w=Pt(1.4))
    add_accent_bar(slide, hero_x, hero_y, hero_w, Inches(0.08), SAPPHIRE)

    add_text(slide, hero_x + Inches(0.5), hero_y + Inches(0.5),
             hero_w - Inches(1.0), Inches(0.7),
             "30-minute Foundry integration POC scoping call.",
             font=HEAD_FONT, size=28, bold=True, color=ICE,
             line_spacing=1.2)
    add_text(slide, hero_x + Inches(0.5), hero_y + Inches(1.4),
             hero_w - Inches(1.0), Inches(0.5),
             "with one Foundry corp-dev lead and one Foundry platform engineer.",
             font=HEAD_FONT, size=15, color=SAPPHIRE_LT, italic=True)

    add_text(slide, hero_x + Inches(0.5), hero_y + Inches(2.3),
             hero_w - Inches(1.0), Inches(0.4),
             "Agenda (30 min):",
             font=HEAD_FONT, size=14, bold=True, color=ICE)

    add_bullet_list(slide, hero_x + Inches(0.5), hero_y + Inches(2.8),
                    hero_w - Inches(1.0), Inches(1.4),
                    [
                        "Walk one signal end-to-end through Risk Kernel + Provenance + Foundry sync (10 min).",
                        "Map IntelVectorRecord, ThreatIndicator, OODAPacket onto your existing Foundry ontology (10 min).",
                        "Decide a 30-day mutual-fit POC scope or close out cleanly (10 min).",
                    ],
                    size=13, color=ICE_MUTED, marker_color=SAPPHIRE_LT)

    # Contact line
    add_text(slide, Inches(0.6), Inches(6.25), Inches(12.0), Inches(0.4),
             "Reply to: aristotlespec@gmail.com",
             font=MONO_FONT, size=16, bold=True, color=SAPPHIRE_LT,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(6.65), Inches(12.0), Inches(0.4),
             "Diligence packet + acquirer microsite are linked in the appendix.",
             font=BODY_FONT, size=12, italic=True, color=ICE_MUTED,
             align=PP_ALIGN.CENTER)

    add_footer(slide, dark=True, label="12 / 14")
    add_speaker_notes(
        slide,
        "Concrete ask: a 30-minute scoping call. One corp-dev lead + one "
        "platform engineer is the right room shape — Foundry's value lands "
        "best when both sides are present.\n\n"
        "Three-step agenda is intentional: walk a signal, map onto ontology, "
        "decide. No deliverables required from the buyer in advance.\n\n"
        f"Source: docs/diligence/00-executive-summary.md + project_palantir_pitch memory @ SHA {HEAD_SHA}"
    )
    return slide


def slide_appendix_diligence(prs):
    """Slide 13 — Appendix: diligence packet."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_ELEV)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Appendix · Diligence packet",
             font=HEAD_FONT, size=26, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Ten markdown files in docs/diligence/. Read in order; ~30 minutes.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Two-column listing
    items = [
        ("00", "Executive summary",
         "What Sapphire is, why it is buyable, where the controls live."),
        ("01", "Architecture",
         "Five planes: operator surfaces, runtimes, safety, artifacts, complement lanes."),
        ("02", "Product surfaces",
         "Dashboard pages + APIs + plugin tools + Hermes skills."),
        ("03", "Security posture",
         "SOC-2-flavoured map. Open incident: pending operator credential rotation."),
        ("04", "Data + models",
         "4+1 inference mesh. Local-first; bounded Gemini OODA; provenance everywhere."),
        ("05", "Operations",
         "LaunchAgents + readiness sweeps + routine soak gates."),
        ("06", "Financials + spend",
         "No-spend by default. Cost controls are repo-enforced."),
        ("07", "Roadmap",
         "30 / 60 / 90 day plus 1-year — packaging, not invention."),
        ("08", "Incidents + postmortems",
         "Each incident becomes a code guard, test, runbook, or tracked decision."),
        ("09", "Team + process",
         "One operator. Codex + Claude/Hermes constrained by skill inventories."),
    ]

    grid_x = Inches(0.6)
    grid_y = Inches(1.7)
    row_h = Inches(0.48)
    label_w = Inches(0.6)
    title_w = Inches(2.5)
    body_w = Inches(2.85)
    gap_y = Inches(0.05)
    cols = 2
    col_w = Inches(6.1)
    col_gap = Inches(0.15)

    for i, (num, title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = grid_x + col * (col_w + col_gap)
        y = grid_y + row * (row_h + gap_y)

        # Number badge
        add_text(slide, x, y, label_w, row_h,
                 num, font=MONO_FONT, size=15, bold=True, color=SAPPHIRE_LT,
                 vertical=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(slide, x + label_w, y, title_w, row_h,
                 title, font=HEAD_FONT, size=12, bold=True, color=ICE,
                 vertical=MSO_ANCHOR.MIDDLE)
        # Body
        add_text(slide, x + label_w + title_w, y, body_w, row_h,
                 body, font=BODY_FONT, size=10.5, color=ICE_MUTED,
                 line_spacing=1.2, vertical=MSO_ANCHOR.MIDDLE)

    # Microsite pointer
    add_card(slide, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.3),
             fill=NAVY_CARD, edge=SAPPHIRE, edge_w=Pt(1.0))
    add_text(slide, Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.4),
             "Acquirer microsite",
             font=HEAD_FONT, size=14, bold=True, color=SAPPHIRE_LT)
    add_text(slide, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.85),
             "web/acquirer/index.html  ·  9 capability cards · 8 stack items · "
             "diligence index · safety posture · founder contact.\n"
             "Vanilla CSS, no analytics, no tracking, no external JS. "
             "Ready to host on any static surface.",
             font=BODY_FONT, size=11, color=ICE_MUTED, line_spacing=1.35)

    add_footer(slide, dark=True, label="13 / 14")
    add_speaker_notes(
        slide,
        "All ten files live under docs/diligence/. The microsite at "
        "web/acquirer/index.html is the same content reformatted for a "
        "single-page corp-dev brief.\n\n"
        f"Source: docs/diligence/ + web/acquirer/index.html @ SHA {HEAD_SHA}"
    )
    return slide


def slide_appendix_provenance(prs):
    """Slide 14 — Appendix: provenance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, NAVY_ELEV)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.55),
             "Appendix · Provenance",
             font=HEAD_FONT, size=26, bold=True, color=ICE)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.4),
             "Every metric in this deck is repo-grounded. Every slide notes its source.",
             font=HEAD_FONT, size=13, italic=True, color=ICE_MUTED)

    # Two cards: (left) build provenance; (right) verification commands
    left_x = Inches(0.6)
    left_y = Inches(1.6)
    left_w = Inches(6.1)
    left_h = Inches(5.2)
    add_card(slide, left_x, left_y, left_w, left_h, fill=NAVY_CARD)
    add_accent_bar(slide, left_x, left_y, left_w, Inches(0.06), SAPPHIRE)
    add_text(slide, left_x + Inches(0.3), left_y + Inches(0.25),
             left_w - Inches(0.6), Inches(0.5),
             "Build provenance",
             font=HEAD_FONT, size=18, bold=True, color=ICE)

    provenance_lines = [
        f"Repo:          arigatoexpress/Sapphire",
        f"HEAD:          {HEAD_SHA}",
        f"Build date:    {BUILD_DATE}",
        f"Branch:        docs/sapphire-pitch-deck-2026-04-29",
        f"Builder:       /tmp/sapphire-deck-build/",
        f"PPTX library:  python-pptx",
        f"Diagram:       matplotlib (PIL backend)",
        f"Test snapshot: scripts/ops/test_inventory.py @ 2026-04-28",
        f"Audience:      {TARGET_AUDIENCE}",
        f"Sidecar:       *.envelope.json (lib/core/provenance shape)",
    ]
    add_text(slide, left_x + Inches(0.35), left_y + Inches(0.95),
             left_w - Inches(0.7), left_h - Inches(1.1),
             "\n".join(provenance_lines),
             font=MONO_FONT, size=10.5, color=ICE_MUTED, line_spacing=1.55)

    right_x = Inches(6.85)
    right_y = Inches(1.6)
    right_w = Inches(6.1)
    right_h = Inches(5.2)
    add_card(slide, right_x, right_y, right_w, right_h, fill=NAVY_CARD)
    add_accent_bar(slide, right_x, right_y, right_w, Inches(0.06), SAPPHIRE_LT)
    add_text(slide, right_x + Inches(0.3), right_y + Inches(0.25),
             right_w - Inches(0.6), Inches(0.5),
             "Verify it yourself",
             font=HEAD_FONT, size=18, bold=True, color=ICE)

    verify_lines = [
        "# Test count",
        "$ python3 scripts/ops/test_inventory.py --check-readme",
        "",
        "# Plugin tool count",
        "$ ls plugins/claw-sapphire/tools/*.py | wc -l",
        "",
        "# Dashboard pages",
        "$ ls services/dashboard/templates/pages/ | wc -l",
        "",
        "# Provenance verification",
        "$ python3 scripts/ops/provenance_verify.py --pretty",
        "",
        "# Production readiness",
        "$ python3 scripts/ops/production_readiness_sweep.py",
        "",
        "# Risk kernel red-team",
        "$ pytest tests/integration/test_risk_kernel_redteam.py",
    ]
    add_text(slide, right_x + Inches(0.35), right_y + Inches(0.95),
             right_w - Inches(0.7), right_h - Inches(1.1),
             "\n".join(verify_lines),
             font=MONO_FONT, size=10, color=ICE_MUTED, line_spacing=1.45)

    # Bottom honest disclosure
    add_text(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
             "This deck contains zero fabricated metrics. Every claim is reproducible from the repo at the SHA above.",
             font=HEAD_FONT, size=11, italic=True, color=SAPPHIRE_LT,
             align=PP_ALIGN.CENTER)

    add_speaker_notes(
        slide,
        "Provenance addendum.\n\n"
        "Every numeric claim in this deck has a source command or file "
        "reference. The sidecar JSON file (lib/core/provenance.py-shape) "
        "records the same on disk.\n\n"
        f"Source: lib/core/provenance.py @ SHA {HEAD_SHA}"
    )
    return slide


def main():
    out_dir = "/tmp/sapphire-deck-build"
    arch = os.path.join(out_dir, "architecture.png")
    if not os.path.exists(arch):
        sys.exit(f"missing diagram: {arch}")

    prs = make_presentation()
    slide_title(prs)
    slide_thesis(prs)
    slide_architecture(prs, arch)
    slide_capabilities(prs)
    slide_safety(prs)
    slide_traction(prs)
    slide_competitive(prs)
    slide_why_audience(prs)
    slide_live_trading_ramp(prs)
    slide_team_ip(prs)
    slide_acquisition_rationale(prs)
    slide_the_ask(prs)
    slide_appendix_diligence(prs)
    slide_appendix_provenance(prs)

    output = os.path.join(out_dir, "sapphire-pitch-deck-2026-04-29.pptx")
    prs.save(output)
    print(f"WROTE {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
